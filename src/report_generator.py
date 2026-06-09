import pandas as pd
import numpy as np
from io import BytesIO
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, HRFlowable
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import warnings
warnings.filterwarnings('ignore')

def generate_excel_report(df, stats_df=None, title="Analytics Report"):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        wb = writer.book
        hdr = wb.add_format({'bold':True,'font_color':'#FFFFFF','bg_color':'#1565C0','align':'center','border':1,'font_size':11})
        col_hdr = wb.add_format({'bold':True,'font_color':'#FFFFFF','bg_color':'#2196F3','align':'center','border':1,'font_size':10})
        alt = wb.add_format({'bg_color':'#E3F2FD','border':1,'font_size':9})
        nrm = wb.add_format({'bg_color':'#FFFFFF','border':1,'font_size':9})
        grn = wb.add_format({'bold':True,'font_color':'#FFFFFF','bg_color':'#2E7D32','align':'center','border':1,'font_size':10})
        org = wb.add_format({'bold':True,'font_color':'#FFFFFF','bg_color':'#E65100','align':'center','border':1,'font_size':10})

        # Sheet 1
        df.to_excel(writer, sheet_name='Raw Data', index=False, startrow=2)
        ws = writer.sheets['Raw Data']
        ws.merge_range(0,0,0,min(len(df.columns)-1,9), f'InsightPro BI — {title}', hdr)
        for i,c in enumerate(df.columns):
            ws.write(2,i,c,col_hdr)
            ws.set_column(i,i,16)
        for r in range(min(len(df),500)):
            for c in range(len(df.columns)):
                try: ws.write(r+3,c,df.iloc[r,c], alt if r%2==0 else nrm)
                except: ws.write(r+3,c,str(df.iloc[r,c]), alt if r%2==0 else nrm)

        # Sheet 2
        df.describe().round(2).to_excel(writer, sheet_name='Summary Statistics', startrow=2)
        ws2 = writer.sheets['Summary Statistics']
        ws2.merge_range(0,0,0,min(len(df.columns),8),'Summary Statistics',grn)
        ws2.set_column(0,10,16)

        # Sheet 3
        comp = round((1-df.isnull().sum().sum()/(len(df)*len(df.columns)))*100,2)
        qdf = pd.DataFrame({
            'Metric':['Total Rows','Total Columns','Missing Values','Duplicate Rows','Numeric Cols','Categorical Cols','Completeness %'],
            'Value':[len(df),len(df.columns),df.isnull().sum().sum(),df.duplicated().sum(),
                    len(df.select_dtypes(include=np.number).columns),
                    len(df.select_dtypes(include='object').columns),f"{comp}%"],
            'Status':['Good','Good',
                     'Clean' if df.isnull().sum().sum()==0 else 'Has Missing',
                     'No Dups' if df.duplicated().sum()==0 else 'Has Dups',
                     'Available','Available','Excellent' if comp>95 else 'Good']
        })
        qdf.to_excel(writer, sheet_name='Data Quality', index=False, startrow=2)
        ws3 = writer.sheets['Data Quality']
        ws3.merge_range(0,0,0,2,'Data Quality Report',org)
        for i,c in enumerate(qdf.columns): ws3.write(2,i,c,col_hdr)
        ws3.set_column(0,2,25)

        # Sheet 4
        miss = df.isnull().sum().to_frame('Missing Count')
        miss['Missing %'] = (df.isnull().sum()/len(df)*100).round(2)
        miss['Status'] = miss['Missing Count'].apply(lambda x:'Complete' if x==0 else('Minor' if x<len(df)*0.1 else 'Major'))
        miss.to_excel(writer, sheet_name='Missing Values', startrow=2)
        ws4 = writer.sheets['Missing Values']
        ws4.merge_range(0,0,0,3,'Missing Values Analysis',org)
        ws4.set_column(0,3,22)

        # Sheet 5
        ncols = df.select_dtypes(include=np.number).columns
        if len(ncols)>0:
            nd = []
            for col in ncols:
                d = df[col].dropna()
                q1,q3 = d.quantile(0.25),d.quantile(0.75)
                out = len(df[(df[col]<q1-1.5*(q3-q1))|(df[col]>q3+1.5*(q3-q1))])
                nd.append({'Column':col,'Mean':round(d.mean(),3),'Median':round(d.median(),3),
                          'Std Dev':round(d.std(),3),'Min':round(d.min(),3),'Max':round(d.max(),3),
                          'Skewness':round(d.skew(),3),'Outliers':out})
            ndf = pd.DataFrame(nd)
            ndf.to_excel(writer, sheet_name='Numeric Analysis', index=False, startrow=2)
            ws5 = writer.sheets['Numeric Analysis']
            ws5.merge_range(0,0,0,7,'Numeric Column Analysis',grn)
            for i,c in enumerate(ndf.columns): ws5.write(2,i,c,col_hdr)
            ws5.set_column(0,7,16)

        # Sheet 6
        if stats_df is not None and len(stats_df)>0:
            stats_df.to_excel(writer, sheet_name='Advanced Stats', startrow=2)
            ws6 = writer.sheets['Advanced Stats']
            ws6.merge_range(0,0,0,min(len(stats_df.columns),10),'Advanced Stats',grn)
            ws6.set_column(0,10,16)
    return output.getvalue()

def generate_pdf_report(df, title="Analytics Report", insights=None):
    output = BytesIO()
    doc = SimpleDocTemplate(output, pagesize=A4,
                           rightMargin=0.6*inch, leftMargin=0.6*inch,
                           topMargin=0.6*inch, bottomMargin=0.6*inch)
    elements = []
    sec = ParagraphStyle('S',fontSize=13,textColor=colors.HexColor('#1565C0'),
                         fontName='Helvetica-Bold',spaceBefore=10,spaceAfter=5)

    # Header
    h = Table([[title]], colWidths=[6.8*inch])
    h.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#1565C0')),
        ('TEXTCOLOR',(0,0),(-1,-1),colors.white),
        ('FONTNAME',(0,0),(-1,-1),'Helvetica-Bold'),
        ('FONTSIZE',(0,0),(-1,-1),18),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('TOPPADDING',(0,0),(-1,-1),14),
        ('BOTTOMPADDING',(0,0),(-1,-1),14),
    ]))
    elements.append(h)

    sub = Table([['Generated by InsightPro BI Platform  |  Krishan Kumar Chauhan  |  M.Tech Data Science, GBU']],
                colWidths=[6.8*inch])
    sub.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#1976D2')),
        ('TEXTCOLOR',(0,0),(-1,-1),colors.HexColor('#BBDEFB')),
        ('FONTNAME',(0,0),(-1,-1),'Helvetica'),
        ('FONTSIZE',(0,0),(-1,-1),8),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('TOPPADDING',(0,0),(-1,-1),5),
        ('BOTTOMPADDING',(0,0),(-1,-1),5),
    ]))
    elements.append(sub)
    elements.append(Spacer(1,0.15*inch))

    # KPI Cards
    elements.append(Paragraph("Dataset Overview", sec))
    comp = round((1-df.isnull().sum().sum()/(len(df)*len(df.columns)))*100,1)
    nc = len(df.select_dtypes(include=np.number).columns)
    cc = len(df.select_dtypes(include='object').columns)

    miss_color = colors.HexColor('#C62828') if df.isnull().sum().sum()>0 else colors.HexColor('#2E7D32')
    dup_color = colors.HexColor('#C62828') if df.duplicated().sum()>0 else colors.HexColor('#2E7D32')

    kpi_headers = [['Total Rows','Columns','Missing','Duplicates','Completeness']]
    kpi_values = [[f"{len(df):,}", str(len(df.columns)),
                   str(df.isnull().sum().sum()), str(df.duplicated().sum()), f"{comp}%"]]

    kpi_h_table = Table(kpi_headers, colWidths=[1.36*inch]*5)
    kpi_h_table.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(0,0),colors.HexColor('#1565C0')),
        ('BACKGROUND',(1,0),(1,0),colors.HexColor('#2E7D32')),
        ('BACKGROUND',(2,0),(2,0),miss_color),
        ('BACKGROUND',(3,0),(3,0),dup_color),
        ('BACKGROUND',(4,0),(4,0),colors.HexColor('#00695C')),
        ('TEXTCOLOR',(0,0),(-1,-1),colors.white),
        ('FONTNAME',(0,0),(-1,-1),'Helvetica-Bold'),
        ('FONTSIZE',(0,0),(-1,-1),8),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('TOPPADDING',(0,0),(-1,-1),8),
        ('BOTTOMPADDING',(0,0),(-1,-1),8),
        ('GRID',(0,0),(-1,-1),1,colors.white),
    ]))
    elements.append(kpi_h_table)

    kpi_v_table = Table(kpi_values, colWidths=[1.36*inch]*5)
    kpi_v_table.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(0,0),colors.HexColor('#E3F2FD')),
        ('BACKGROUND',(1,0),(1,0),colors.HexColor('#E8F5E9')),
        ('BACKGROUND',(2,0),(2,0),colors.HexColor('#FFEBEE') if df.isnull().sum().sum()>0 else colors.HexColor('#E8F5E9')),
        ('BACKGROUND',(3,0),(3,0),colors.HexColor('#FFEBEE') if df.duplicated().sum()>0 else colors.HexColor('#E8F5E9')),
        ('BACKGROUND',(4,0),(4,0),colors.HexColor('#E0F2F1')),
        ('TEXTCOLOR',(0,0),(0,0),colors.HexColor('#1565C0')),
        ('TEXTCOLOR',(1,0),(1,0),colors.HexColor('#2E7D32')),
        ('TEXTCOLOR',(2,0),(2,0),colors.HexColor('#C62828') if df.isnull().sum().sum()>0 else colors.HexColor('#2E7D32')),
        ('TEXTCOLOR',(3,0),(3,0),colors.HexColor('#C62828') if df.duplicated().sum()>0 else colors.HexColor('#2E7D32')),
        ('TEXTCOLOR',(4,0),(4,0),colors.HexColor('#00695C')),
        ('FONTNAME',(0,0),(-1,-1),'Helvetica-Bold'),
        ('FONTSIZE',(0,0),(-1,-1),16),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('TOPPADDING',(0,0),(-1,-1),10),
        ('BOTTOMPADDING',(0,0),(-1,-1),10),
        ('GRID',(0,0),(-1,-1),1,colors.white),
    ]))
    elements.append(kpi_v_table)
    elements.append(Spacer(1,0.15*inch))

    # Summary Statistics
    num_cols = df.select_dtypes(include=np.number).columns[:5]
    if len(num_cols)>0:
        elements.append(Paragraph("Summary Statistics", sec))
        stats = df[num_cols].describe().round(2)
        sd = [['Statistic']+list(num_cols)]
        for idx in stats.index:
            sd.append([idx.capitalize()]+[str(stats.loc[idx,c]) for c in num_cols])
        cw = [0.9*inch]+[5.9*inch/len(num_cols)]*len(num_cols)
        st = Table(sd, colWidths=cw)
        st.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#1565C0')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.white),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,-1),8),
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#BBDEFB')),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.HexColor('#E3F2FD'),colors.white]),
            ('TOPPADDING',(0,0),(-1,-1),5),
            ('BOTTOMPADDING',(0,0),(-1,-1),5),
        ]))
        elements.append(st)
        elements.append(Spacer(1,0.15*inch))

    # Numeric Analysis
    if len(num_cols)>0:
        elements.append(Paragraph("Numeric Column Analysis", sec))
        na = [['Column','Mean','Median','Std Dev','Min','Max','Skewness','Outliers']]
        for col in num_cols:
            d = df[col].dropna()
            q1,q3 = d.quantile(0.25),d.quantile(0.75)
            out = len(df[(df[col]<q1-1.5*(q3-q1))|(df[col]>q3+1.5*(q3-q1))])
            na.append([col,f"{d.mean():.2f}",f"{d.median():.2f}",
                      f"{d.std():.2f}",f"{d.min():.2f}",f"{d.max():.2f}",
                      f"{d.skew():.3f}",str(out)])
        cw2 = [1.1*inch]+[0.83*inch]*7
        nt = Table(na, colWidths=cw2)
        nt.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),colors.HexColor('#2E7D32')),
            ('TEXTCOLOR',(0,0),(-1,0),colors.white),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,-1),8),
            ('ALIGN',(0,0),(-1,-1),'CENTER'),
            ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#C8E6C9')),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.HexColor('#E8F5E9'),colors.white]),
            ('TOPPADDING',(0,0),(-1,-1),5),
            ('BOTTOMPADDING',(0,0),(-1,-1),5),
        ]))
        elements.append(nt)
        elements.append(Spacer(1,0.15*inch))

    # Insights
    elements.append(HRFlowable(width="100%",thickness=1.5,color=colors.HexColor('#2196F3')))
    elements.append(Paragraph("Key Insights & Recommendations", sec))

    auto = []
    auto.append("Dataset is clean — no missing values" if df.isnull().sum().sum()==0
                else f"{df.isnull().sum().sum()} missing values — consider imputation")
    auto.append("No duplicate rows — data integrity maintained" if df.duplicated().sum()==0
                else f"{df.duplicated().sum()} duplicates found — review data")
    auto.append(f"Dataset completeness: {comp}% — {'Excellent' if comp>95 else 'Good'}")
    auto.append(f"{nc} numeric + {cc} categorical columns available for analysis")
    auto.append(f"{len(df):,} total records analyzed by InsightPro BI Platform")
    if insights: auto.extend(insights)

    idata = [[Paragraph(f"  {i}", ParagraphStyle('ip',fontSize=9,
              textColor=colors.HexColor('#1B5E20'),fontName='Helvetica',spaceAfter=2))]
             for i in auto]
    it = Table(idata, colWidths=[6.8*inch])
    it.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#F1F8E9')),
        ('LINEBELOW',(0,0),(-1,-2),0.3,colors.HexColor('#C5E1A5')),
        ('LEFTPADDING',(0,0),(-1,-1),10),
        ('TOPPADDING',(0,0),(-1,-1),6),
        ('BOTTOMPADDING',(0,0),(-1,-1),6),
    ]))
    elements.append(it)
    elements.append(Spacer(1,0.1*inch))

    # Footer
    elements.append(HRFlowable(width="100%",thickness=1,color=colors.HexColor('#2196F3')))
    ft = Table([['InsightPro BI Platform  |  Krishan Kumar Chauhan  |  M.Tech Data Science, GBU  |  Python • Pandas • Plotly • Streamlit']],
               colWidths=[6.8*inch])
    ft.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#E3F2FD')),
        ('TEXTCOLOR',(0,0),(-1,-1),colors.HexColor('#1565C0')),
        ('FONTNAME',(0,0),(-1,-1),'Helvetica'),
        ('FONTSIZE',(0,0),(-1,-1),7),
        ('ALIGN',(0,0),(-1,-1),'CENTER'),
        ('TOPPADDING',(0,0),(-1,-1),5),
        ('BOTTOMPADDING',(0,0),(-1,-1),5),
    ]))
    elements.append(ft)
    doc.build(elements)
    return output.getvalue()

def generate_ppt_report(df, title="Analytics Report", charts=None):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def bg(slide, hex_c):
        f = slide.background.fill
        f.solid()
        f.fore_color.rgb = RGBColor.from_string(hex_c)

    def box(slide, l, t, w, h, hex_c):
        s = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor.from_string(hex_c)
        s.line.fill.background()
        return s

    def txt(slide, text, l, t, w, h, sz=14, bold=False, clr='FFFFFF', align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(clr)

    comp = round((1-df.isnull().sum().sum()/(len(df)*len(df.columns)))*100,1)
    nc = len(df.select_dtypes(include=np.number).columns)
    cc = len(df.select_dtypes(include='object').columns)

    # Slide 1 - Title
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s1,'0D1117')
    box(s1,0,0,13.33,0.12,'2196F3')
    box(s1,0,7.38,13.33,0.12,'2196F3')
    box(s1,1.0,1.5,11.33,3.0,'1565C0')
    txt(s1,'InsightPro BI Platform',1.2,1.65,11,0.9,sz=30,bold=True,clr='FFFFFF',align=PP_ALIGN.CENTER)
    txt(s1,title,1.2,2.7,11,0.7,sz=20,clr='90CAF9',align=PP_ALIGN.CENTER)
    txt(s1,'Krishan Kumar Chauhan  |  M.Tech Data Science, GBU',1,4.2,11.33,0.5,sz=13,clr='B0BEC5',align=PP_ALIGN.CENTER)
    txt(s1,f'{len(df):,} rows  x  {len(df.columns)} columns  |  {comp}% Complete',1,4.9,11.33,0.5,sz=12,clr='A5D6A7',align=PP_ALIGN.CENTER)

    # Slide 2 - Overview
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s2,'0A0E1A')
    box(s2,0,0,13.33,0.9,'1565C0')
    txt(s2,'  Dataset Overview',0,0.15,13,0.6,sz=22,bold=True,clr='FFFFFF')
    kpis = [
        ('Total Rows',f"{len(df):,}",'1565C0'),
        ('Columns',str(len(df.columns)),'2E7D32'),
        ('Missing',str(df.isnull().sum().sum()),'C62828' if df.isnull().sum().sum()>0 else '2E7D32'),
        ('Duplicates',str(df.duplicated().sum()),'C62828' if df.duplicated().sum()>0 else '2E7D32'),
        ('Completeness',f"{comp}%",'00695C'),
    ]
    for i,(label,value,color) in enumerate(kpis):
        x = 0.4+i*2.5
        box(s2,x,1.1,2.2,2.0,color)
        txt(s2,value,x,1.25,2.2,0.9,sz=28,bold=True,clr='FFFFFF',align=PP_ALIGN.CENTER)
        txt(s2,label,x,2.2,2.2,0.5,sz=11,clr='E3F2FD',align=PP_ALIGN.CENTER)
    txt(s2,f'Numeric: {nc}   Categorical: {cc}',0.5,3.5,12,0.5,sz=14,clr='90CAF9')
    txt(s2,'Columns: '+'  |  '.join(list(df.columns)[:7]),0.5,4.1,12,0.5,sz=10,clr='78909C')

    # Slide 3 - Statistics
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s3,'0A0E1A')
    box(s3,0,0,13.33,0.9,'2E7D32')
    txt(s3,'  Statistical Summary',0,0.15,13,0.6,sz=22,bold=True,clr='FFFFFF')
    num_cols = df.select_dtypes(include=np.number).columns[:4]
    clrs = ['1565C0','2E7D32','E65100','6A1B9A']
    for i,col in enumerate(num_cols):
        d = df[col].dropna()
        y = 1.1+i*1.5
        box(s3,0.3,y,12.5,1.2,clrs[i%4])
        txt(s3,col,0.5,y+0.1,3,0.4,sz=13,bold=True,clr='FFFFFF')
        txt(s3,f"Mean: {d.mean():.2f}   Median: {d.median():.2f}   Std: {d.std():.2f}   Min: {d.min():.2f}   Max: {d.max():.2f}   Skew: {d.skew():.3f}",
            0.5,y+0.6,12,0.5,sz=11,clr='E3F2FD')

    # Slide 4 - Quality
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s4,'0A0E1A')
    box(s4,0,0,13.33,0.9,'E65100')
    txt(s4,'  Data Quality Report',0,0.15,13,0.6,sz=22,bold=True,clr='FFFFFF')
    quality = [
        ('Missing Values',df.isnull().sum().sum(),'Clean' if df.isnull().sum().sum()==0 else 'Has Missing','2E7D32' if df.isnull().sum().sum()==0 else 'C62828'),
        ('Duplicate Rows',df.duplicated().sum(),'No Dups' if df.duplicated().sum()==0 else 'Has Dups','2E7D32' if df.duplicated().sum()==0 else 'C62828'),
        ('Completeness',f"{comp}%",'Excellent' if comp>95 else 'Good','2E7D32' if comp>95 else 'E65100'),
        ('Numeric Cols',nc,'Available','1565C0'),
        ('Categorical',cc,'Available','6A1B9A'),
    ]
    for i,(metric,value,status,color) in enumerate(quality):
        x = 0.35+(i%3)*4.3
        y = 1.1+(i//3)*2.5
        box(s4,x,y,3.9,2.1,color)
        txt(s4,str(value),x,y+0.2,3.9,0.8,sz=30,bold=True,clr='FFFFFF',align=PP_ALIGN.CENTER)
        txt(s4,metric,x,y+1.1,3.9,0.4,sz=11,clr='E3F2FD',align=PP_ALIGN.CENTER)
        txt(s4,status,x,y+1.55,3.9,0.4,sz=10,clr='FFEB3B',align=PP_ALIGN.CENTER)

    # Slide 5 - Insights
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s5,'0A0E1A')
    box(s5,0,0,13.33,0.9,'6A1B9A')
    txt(s5,'  Key Insights',0,0.15,13,0.6,sz=22,bold=True,clr='FFFFFF')
    ilist = [
        f"Dataset: {len(df):,} rows x {len(df.columns)} columns — {comp}% complete",
        f"Missing: {df.isnull().sum().sum()} — {'Clean & ready for analysis' if df.isnull().sum().sum()==0 else 'Consider imputation'}",
        f"Duplicates: {df.duplicated().sum()} — {'Integrity maintained' if df.duplicated().sum()==0 else 'Review records'}",
        f"Numeric Columns: {nc} — available for ML modeling",
        f"Categorical Columns: {cc} — available for encoding",
        f"Quality: {'Excellent' if comp>95 else 'Good'} ({comp}% completeness)",
    ]
    box_clrs = ['1A237E','1B5E20','0D47A1','311B92','004D40','4A148C']
    for i,insight in enumerate(ilist):
        y = 1.1+i*0.95
        box(s5,0.3,y,12.5,0.8,box_clrs[i])
        txt(s5,f"  {insight}",0.4,y+0.1,12,0.6,sz=12,clr='E3F2FD')

    # Slide 6 - Thank You
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s6,'0D1117')
    box(s6,0,0,13.33,0.12,'2196F3')
    box(s6,0,7.38,13.33,0.12,'2196F3')
    box(s6,1.5,2.0,10.33,3.5,'1565C0')
    txt(s6,'Thank You!',1.5,2.2,10.33,1.0,sz=36,bold=True,clr='FFFFFF',align=PP_ALIGN.CENTER)
    txt(s6,'InsightPro BI Platform',1.5,3.3,10.33,0.6,sz=18,bold=True,clr='90CAF9',align=PP_ALIGN.CENTER)
    txt(s6,'Krishan Kumar Chauhan  |  M.Tech Data Science, GBU',1.5,4.0,10.33,0.5,sz=13,clr='B0BEC5',align=PP_ALIGN.CENTER)
    txt(s6,'Python  |  Pandas  |  Plotly  |  Streamlit  |  MongoDB  |  SQL  |  Scikit-learn',1.5,4.6,10.33,0.5,sz=11,clr='A5D6A7',align=PP_ALIGN.CENTER)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()